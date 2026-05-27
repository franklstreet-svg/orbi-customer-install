"""
pptx_gen — Tier-2 feature: generate a real PowerPoint deck from a topic.

Pipeline mirrors doc_convert / image_gen:

    topic (string)
      → build_outline()  -- LLM turns the topic into a structured outline
      → render_deck()    -- python-pptx renders the outline into real .pptx bytes
      → save_to_workspace()  -- drop the file into ~/Orbi/ so it shows up in Files

If the LLM is unreachable, build_outline returns a sensible one-slide fallback
so the owner ALWAYS gets a deck back, just like image_gen always returns an
image. The owner is never stuck.

Themes ("modern", "light", "bold") are real visual themes, not just font
swaps — they each set background, text color, and accent-bar color.

CONVENTIONS
-----------
* logging.getLogger("orbi.pptx_gen")
* python-pptx lazy-imported inside functions (other modules follow the same
  pattern; keeps import time low if the deck builder is never used)
* All output is raw .pptx bytes
* 16:9 aspect ratio (13.333" x 7.5") — modern widescreen
* Minimum legible sizing: 24pt body, 40pt headings
"""

from __future__ import annotations

import hashlib
import io
import json
import logging
import re
import time
from datetime import datetime
from pathlib import Path

log = logging.getLogger("orbi.pptx_gen")


# ---------------------------------------------------------------------------
# Brand palette — pulled from owner_dashboard/dashboard.css so the "modern"
# theme matches the rest of the Orbi UI.
# ---------------------------------------------------------------------------

# Background (matches dashboard bg + theme-color meta tag)
MODERN_BG          = (0x0b, 0x0f, 0x1a)   # #0b0f1a — deep slate
MODERN_BG_SOFT     = (0x13, 0x1a, 0x2e)   # #131a2e — subtle band
MODERN_PANEL       = (0x1a, 0x22, 0x36)   # #1a2236 — card background
MODERN_BORDER      = (0x2c, 0x39, 0x57)   # #2c3957 — divider line
MODERN_TEXT        = (0xea, 0xf0, 0xff)   # #eaf0ff — primary white-ish
MODERN_TEXT_MUTED  = (0xb8, 0xc6, 0xe0)   # #b8c6e0 — secondary
MODERN_ACCENT_BLUE = (0x4f, 0x8c, 0xff)   # #4f8cff
MODERN_ACCENT_PURP = (0x8b, 0x5c, 0xf6)   # #8b5cf6

LIGHT_BG    = (0xff, 0xff, 0xff)
LIGHT_TEXT  = (0x1a, 0x22, 0x36)
LIGHT_MUTED = (0x55, 0x66, 0x88)
LIGHT_ACCENT = (0x4f, 0x8c, 0xff)

BOLD_BG     = (0x8b, 0x5c, 0xf6)         # full-bleed purple
BOLD_BG_ALT = (0x4f, 0x8c, 0xff)         # alternate slide bleed
BOLD_TEXT   = (0xff, 0xff, 0xff)
BOLD_ACCENT = (0x0b, 0x0f, 0x1a)

# Deck dimensions — 16:9 widescreen
SLIDE_WIDTH_IN  = 13.333
SLIDE_HEIGHT_IN = 7.5

# Title bar height for "modern" header accent band
TITLE_BAR_HEIGHT_IN = 0.55

# Max content guards
MAX_TOPIC_CHARS   = 500
MAX_OUTLINE_CHARS = 8_000
MAX_SLIDES        = 20      # hard ceiling; LLM ceiling is target_slide_count


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def build_outline(config: dict, topic: str, target_slide_count: int = 7) -> dict:
    """
    Turn a topic into a structured slide outline using the LLM.

    Returns:
        {
            "title": "...",
            "subtitle": "...",
            "slides": [
                {"heading": "...", "bullets": ["...", "..."]},
                ...
            ]
        }

    Defensive — tries json.loads first, falls back to a heading-and-bullets
    text parser. If the LLM is unreachable / errors / returns garbage, returns
    a one-slide fallback so render_deck() always has something to work with.
    """
    topic = (topic or "").strip()
    if not topic:
        raise ValueError("topic is required")
    if len(topic) > MAX_TOPIC_CHARS:
        topic = topic[:MAX_TOPIC_CHARS]

    target_slide_count = max(3, min(MAX_SLIDES, int(target_slide_count or 7)))

    system = _outline_system_prompt(target_slide_count)
    user = (f"Topic: {topic}\n\n"
            f"Produce a {target_slide_count}-slide outline in the JSON format "
            "described above. Output ONLY the JSON object — no preamble.")

    raw = ""
    try:
        import llm_client
        resp = llm_client.generate(config or {}, system,
                                   [{"role": "user", "content": user}])
        raw = (resp.text or "").strip()
        log.info("outline llm tier=%s latency=%dms len=%d",
                 getattr(resp, "tier", "?"),
                 getattr(resp, "latency_ms", 0),
                 len(raw))
    except Exception as exc:    # noqa: BLE001 — fall through
        log.warning("llm unreachable for outline: %s", exc)

    if raw:
        parsed = _parse_outline_json(raw) or _parse_outline_text(raw, topic)
        if parsed and parsed.get("slides"):
            return _normalize_outline(parsed, topic, target_slide_count)

    # Total fallback — single-slide deck so something still renders
    log.info("outline fallback engaged for topic=%r", topic[:60])
    return _fallback_outline(topic)


def render_deck(outline: dict, *, theme: str = "modern",
                business_info: dict | None = None) -> bytes:
    """
    Render an outline dict into a real .pptx file (returned as bytes).

    Args:
        outline:       result of build_outline() (or owner-edited equivalent)
        theme:         "modern" (Orbi brand dark), "light" (white), or "bold"
                       (full-bleed color block per slide)
        business_info: optional {name, contact: {phone, email, website}} —
                       if contact has any values, a "thank you" / contact
                       slide is appended.

    Returns: raw .pptx bytes.
    """
    from pptx import Presentation
    from pptx.util import Inches

    theme = (theme or "modern").lower().strip()
    if theme not in ("modern", "light", "bold"):
        log.info("unknown theme=%r, defaulting to modern", theme)
        theme = "modern"

    outline = _normalize_outline(outline or {}, "Untitled Deck", 7)

    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)

    # Blank layout — we draw everything ourselves so theming is consistent.
    # Layout 6 is "Blank" in the default template.
    blank_layout = prs.slide_layouts[6]

    # ── Title slide ───────────────────────────────────────────────────────
    _add_title_slide(prs, blank_layout, outline.get("title", "Untitled"),
                     outline.get("subtitle", ""), theme=theme)

    # ── Content slides ────────────────────────────────────────────────────
    slides = outline.get("slides", [])[:MAX_SLIDES]
    for i, slide_def in enumerate(slides):
        _add_content_slide(
            prs, blank_layout,
            heading=slide_def.get("heading", ""),
            bullets=slide_def.get("bullets", []),
            theme=theme,
            slide_index=i,
        )

    # ── Final contact / thank-you slide (only if contact info provided) ──
    if business_info:
        contact = (business_info.get("contact") or {}) if isinstance(business_info, dict) else {}
        has_contact = any((contact.get(k) or "").strip()
                          for k in ("phone", "email", "website", "address"))
        if has_contact:
            _add_contact_slide(prs, blank_layout, business_info, theme=theme)

    buf = io.BytesIO()
    prs.save(buf)
    data = buf.getvalue()
    log.info("rendered deck: theme=%s slides=%d bytes=%d",
             theme, len(slides) + 1 + (1 if business_info else 0), len(data))
    return data


def build_deck(config: dict, topic: str, target_slide_count: int = 7,
               theme: str = "modern",
               business_info: dict | None = None) -> dict:
    """
    End-to-end. Returns {pptx_bytes, outline, slide_count, theme}.
    """
    outline = build_outline(config or {}, topic, target_slide_count)
    pptx_bytes = render_deck(outline, theme=theme, business_info=business_info)
    # slide_count = title + content slides + optional contact slide
    slide_count = 1 + len(outline.get("slides", []))
    if business_info and any((business_info.get("contact") or {}).get(k)
                              for k in ("phone", "email", "website", "address")):
        slide_count += 1
    return {
        "pptx_bytes":  pptx_bytes,
        "outline":     outline,
        "slide_count": slide_count,
        "theme":       theme,
        "bytes":       len(pptx_bytes),
    }


def save_deck_to_workspace(pptx_bytes: bytes, topic: str,
                           workspace_dir: Path) -> Path:
    """
    Drop the .pptx into the owner's workspace under /decks/ with a stable,
    sortable filename. Matches image_gen.save_to_workspace pattern.

    Filename: pitch_deck_<slug>_<YYYY-MM-DD>_<HHMMSS>.pptx
    """
    workspace_dir = Path(workspace_dir).expanduser()
    decks_dir = workspace_dir / "decks"
    decks_dir.mkdir(parents=True, exist_ok=True)

    now = datetime.now()
    slug = _slugify(topic) or "deck"
    fname = (f"pitch_deck_{slug}"
             f"_{now.strftime('%Y-%m-%d')}"
             f"_{now.strftime('%H%M%S')}.pptx")
    path = decks_dir / fname

    # Atomic-ish write
    tmp = path.with_suffix(".pptx.tmp")
    tmp.write_bytes(pptx_bytes)
    tmp.replace(path)

    log.info("saved deck (%d bytes) -> %s", len(pptx_bytes), path)
    return path


# ---------------------------------------------------------------------------
# Outline parsing & normalization
# ---------------------------------------------------------------------------


def _outline_system_prompt(target_slide_count: int) -> str:
    return (
        "You are a presentation-deck outline assistant. Given a topic, you "
        "produce a clean, punchy slide outline. Output STRICT JSON only, in "
        "this exact shape:\n\n"
        "{\n"
        '  "title": "the deck title (short, max 8 words)",\n'
        '  "subtitle": "one supporting line, max 12 words",\n'
        '  "slides": [\n'
        '    {"heading": "slide heading (4-6 words)",\n'
        '     "bullets": ["bullet 1", "bullet 2", "bullet 3"]},\n'
        "    ...\n"
        "  ]\n"
        "}\n\n"
        "Rules:\n"
        f"1. Produce EXACTLY {target_slide_count} entries in `slides`.\n"
        "2. Each slide gets 2-5 bullets. Each bullet is one short sentence "
        "or fragment, max 14 words.\n"
        "3. No markdown, no asterisks, no numbering inside bullets.\n"
        "4. The final slide should be a call-to-action (reserve, book, call, "
        "visit, sign up — whatever matches the topic).\n"
        "5. Use the owner's language — concrete, specific, no marketing fluff "
        "like \"unlock\" / \"leverage\" / \"synergy\".\n"
        "6. Output ONLY the JSON. No preamble, no code fences, no trailing text."
    )


def _parse_outline_json(raw: str) -> dict | None:
    """Try to extract a JSON object from the LLM's response. Handles common
    failure modes: code fences, leading 'Here is...' text, trailing commentary."""
    s = raw.strip()

    # Strip code fences if the LLM ignored the no-backticks rule
    fence = re.match(r"^```(?:json)?\s*(.+?)\s*```\s*$", s, re.DOTALL)
    if fence:
        s = fence.group(1).strip()

    # Find first { ... last }
    start = s.find("{")
    end   = s.rfind("}")
    if start == -1 or end == -1 or end <= start:
        return None
    candidate = s[start:end + 1]
    try:
        obj = json.loads(candidate)
    except json.JSONDecodeError:
        # Try to repair common mistakes — trailing commas
        repaired = re.sub(r",\s*([}\]])", r"\1", candidate)
        try:
            obj = json.loads(repaired)
        except json.JSONDecodeError:
            return None
    if not isinstance(obj, dict) or "slides" not in obj:
        return None
    return obj


def _parse_outline_text(raw: str, topic: str) -> dict | None:
    """Fallback parser — heading + bullets in plain text. Format like:

        Title: Memorial Day Weekend
        Subtitle: May 24-26

        # What's Happening
        - bullet 1
        - bullet 2

        # The Menu
        - bullet
    """
    if not raw or not raw.strip():
        return None

    title    = ""
    subtitle = ""
    slides: list[dict] = []
    current: dict | None = None

    for line in raw.splitlines():
        line = line.rstrip()
        if not line.strip():
            continue
        low = line.strip().lower()

        # Title: / Subtitle:
        if low.startswith("title:"):
            title = line.split(":", 1)[1].strip().strip('"').strip("'")
            continue
        if low.startswith("subtitle:"):
            subtitle = line.split(":", 1)[1].strip().strip('"').strip("'")
            continue

        # Heading: leading #, ##, "Slide N:", or just a non-bullet line in caps
        m_heading = re.match(r"^\s*(?:#{1,3}\s+|slide\s*\d+\s*[:.\-]?\s*)(.+)$",
                             line, re.IGNORECASE)
        if m_heading and not line.lstrip().startswith(("-", "*", "•")):
            if current and current.get("bullets"):
                slides.append(current)
            current = {"heading": m_heading.group(1).strip().strip('"'),
                       "bullets": []}
            continue

        # Bullet: - foo / * foo / • foo / 1. foo
        m_bullet = re.match(r"^\s*(?:[-*•]|\d+[.)])\s+(.+)$", line)
        if m_bullet:
            if current is None:
                current = {"heading": "Overview", "bullets": []}
            current["bullets"].append(m_bullet.group(1).strip())
            continue

        # Plain text after a heading → treat as a bullet
        if current is not None:
            current["bullets"].append(line.strip())

    if current and current.get("bullets"):
        slides.append(current)

    if not slides:
        return None

    return {
        "title":    title or topic[:60],
        "subtitle": subtitle,
        "slides":   slides,
    }


def _normalize_outline(outline: dict, topic_fallback: str,
                       target_slide_count: int) -> dict:
    """Ensure required keys exist and values are sane types. Trims long fields,
    drops empty slides, hard-caps slide count."""
    out = dict(outline or {})

    title = str(out.get("title") or topic_fallback or "Untitled").strip()
    if len(title) > 120:
        title = title[:117] + "..."

    subtitle = str(out.get("subtitle") or "").strip()
    if len(subtitle) > 200:
        subtitle = subtitle[:197] + "..."

    slides_in = out.get("slides") or []
    if not isinstance(slides_in, list):
        slides_in = []

    slides_out: list[dict] = []
    for s in slides_in:
        if not isinstance(s, dict):
            continue
        heading = str(s.get("heading") or "").strip()
        bullets_raw = s.get("bullets") or []
        if not isinstance(bullets_raw, list):
            bullets_raw = [str(bullets_raw)]
        bullets = []
        for b in bullets_raw:
            t = str(b).strip().lstrip("-*• ").strip()
            if not t:
                continue
            if len(t) > 240:
                t = t[:237] + "..."
            bullets.append(t)
        if heading or bullets:
            if not heading:
                heading = "Details"
            slides_out.append({"heading": heading, "bullets": bullets[:6]})

    slides_out = slides_out[:max(target_slide_count, MAX_SLIDES)]

    return {"title": title, "subtitle": subtitle, "slides": slides_out}


def _fallback_outline(topic: str) -> dict:
    """When the LLM is offline, produce something coherent rather than nothing.
    A single slide that surfaces the topic itself so the deck still opens."""
    short = topic.strip()
    if len(short) > 80:
        short = short[:77] + "..."
    return {
        "title":    short or "Your Presentation",
        "subtitle": "Generated offline — edit this deck to add your own content.",
        "slides": [
            {
                "heading": "About This Deck",
                "bullets": [
                    f"Topic: {short or '(blank)'}",
                    "The AI assistant was offline when this deck was built.",
                    "Open the file in PowerPoint or Google Slides to edit.",
                    "Try again later for a full AI-generated outline.",
                ],
            },
        ],
    }


# ---------------------------------------------------------------------------
# Slide rendering helpers (python-pptx)
# ---------------------------------------------------------------------------


def _rgb(rgb_tuple):
    from pptx.dml.color import RGBColor
    return RGBColor(rgb_tuple[0], rgb_tuple[1], rgb_tuple[2])


def _theme_colors(theme: str, slide_index: int = 0) -> dict:
    """Return the color palette for the given theme. For 'bold', alternate
    between two backgrounds on consecutive slides for visual variety."""
    if theme == "light":
        return {
            "bg":      LIGHT_BG,
            "text":    LIGHT_TEXT,
            "muted":   LIGHT_MUTED,
            "accent":  LIGHT_ACCENT,
            "accent2": MODERN_ACCENT_PURP,
            "bar":     LIGHT_ACCENT,
        }
    if theme == "bold":
        bg = BOLD_BG if (slide_index % 2 == 0) else BOLD_BG_ALT
        return {
            "bg":      bg,
            "text":    BOLD_TEXT,
            "muted":   (235, 235, 245),
            "accent":  BOLD_ACCENT,
            "accent2": BOLD_ACCENT,
            "bar":     BOLD_ACCENT,
        }
    # modern (default)
    return {
        "bg":      MODERN_BG,
        "text":    MODERN_TEXT,
        "muted":   MODERN_TEXT_MUTED,
        "accent":  MODERN_ACCENT_BLUE,
        "accent2": MODERN_ACCENT_PURP,
        "bar":     MODERN_ACCENT_PURP,
    }


def _paint_background(slide, color_rgb) -> None:
    """python-pptx doesn't expose slide background fill directly on every
    layout, so we draw a full-bleed rectangle. Works on every theme."""
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(SLIDE_WIDTH_IN), Inches(SLIDE_HEIGHT_IN),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color_rgb)
    shape.line.fill.background()    # no border
    # Send to back so other shapes sit on top
    try:
        spTree = shape._element.getparent()
        spTree.remove(shape._element)
        spTree.insert(2, shape._element)
    except Exception:    # noqa: BLE001
        pass


def _add_title_bar(slide, colors: dict) -> None:
    """The accent bar at the top of content slides — the "Orbi" signature
    detail in the 'modern' theme. On light/bold themes it's a thin underline
    accent instead."""
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(SLIDE_WIDTH_IN), Inches(TITLE_BAR_HEIGHT_IN),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(colors["bar"])
    bar.line.fill.background()


def _add_orb_accent(slide, colors: dict) -> None:
    """Small circular "orb" badge in the top-left of the title slide — matches
    Orbi's brand mark."""
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE

    orb_size = 0.5
    margin   = 0.6
    orb = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(margin), Inches(margin),
        Inches(orb_size), Inches(orb_size),
    )
    orb.fill.solid()
    orb.fill.fore_color.rgb = _rgb(colors["accent2"])
    orb.line.color.rgb = _rgb(colors["accent"])
    try:
        orb.line.width = _emu_pt(2)
    except Exception:    # noqa: BLE001
        pass


def _emu_pt(pt: float):
    from pptx.util import Pt
    return Pt(pt)


def _add_text_box(slide, left_in: float, top_in: float,
                  width_in: float, height_in: float,
                  text: str, *, font_size_pt: int,
                  color_rgb,
                  bold: bool = False,
                  align: str = "left") -> None:
    """Wrapper around python-pptx text boxes that ALWAYS sets up a paragraph
    + run explicitly. python-pptx's default text frame has no run, so direct
    .text= sets work but the formatting silently drops. We use explicit runs
    every time."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    box = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in),
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)

    p = tf.paragraphs[0]
    p.alignment = {
        "left":   PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right":  PP_ALIGN.RIGHT,
    }.get(align, PP_ALIGN.LEFT)

    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_rgb)
    run.font.name = "Calibri"    # widely available sans-serif default
    return box


def _add_bullet_block(slide, left_in: float, top_in: float,
                      width_in: float, height_in: float,
                      bullets: list[str],
                      *, font_size_pt: int,
                      color_rgb,
                      accent_rgb) -> None:
    """Render bullets one paragraph per item with a colored bullet char and a
    consistent line spacing. python-pptx's bullet-char support varies by
    PowerPoint version — we use a literal "•" character which renders in
    PowerPoint, Keynote, LibreOffice, and Google Slides reliably."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    box = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in),
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)

    for i, bullet in enumerate(bullets or []):
        # Reuse the first paragraph for bullet 0; add new paragraphs after
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # Space after — keeps bullets from cramming together
        try:
            p.space_after = Pt(8)
        except Exception:    # noqa: BLE001
            pass

        # Bullet marker (colored)
        marker = p.add_run()
        marker.text = "•  "
        marker.font.size = Pt(font_size_pt)
        marker.font.bold = True
        marker.font.color.rgb = _rgb(accent_rgb)
        marker.font.name = "Calibri"

        # Body text
        body = p.add_run()
        body.text = bullet
        body.font.size = Pt(font_size_pt)
        body.font.color.rgb = _rgb(color_rgb)
        body.font.name = "Calibri"


def _add_title_slide(prs, layout, title: str, subtitle: str,
                     *, theme: str) -> None:
    """Centered title + subtitle. 'Modern' theme gets an orb badge accent."""
    colors = _theme_colors(theme, slide_index=0)
    slide = prs.slides.add_slide(layout)

    _paint_background(slide, colors["bg"])

    if theme == "modern":
        _add_orb_accent(slide, colors)

    # Centered title — vertical position ~38% from top
    title_w = SLIDE_WIDTH_IN - 1.5
    _add_text_box(
        slide,
        left_in=0.75, top_in=SLIDE_HEIGHT_IN * 0.34,
        width_in=title_w, height_in=1.5,
        text=title or "Untitled",
        font_size_pt=54,
        color_rgb=colors["text"],
        bold=True,
        align="center",
    )

    # Subtitle below
    if subtitle:
        _add_text_box(
            slide,
            left_in=0.75, top_in=SLIDE_HEIGHT_IN * 0.34 + 1.55,
            width_in=title_w, height_in=1.0,
            text=subtitle,
            font_size_pt=24,
            color_rgb=colors["muted"],
            align="center",
        )

    # Small "Made with Orbi" footer — matches image_gen brand tag
    _add_text_box(
        slide,
        left_in=SLIDE_WIDTH_IN - 2.3, top_in=SLIDE_HEIGHT_IN - 0.45,
        width_in=2.0, height_in=0.35,
        text="Made with Orbi",
        font_size_pt=10,
        color_rgb=colors["muted"],
        align="right",
    )


def _add_content_slide(prs, layout, heading: str, bullets: list[str],
                       *, theme: str, slide_index: int) -> None:
    """Heading at top, bullets in the body. 'Modern' draws a thin accent bar
    along the top edge."""
    colors = _theme_colors(theme, slide_index=slide_index + 1)
    slide = prs.slides.add_slide(layout)

    _paint_background(slide, colors["bg"])

    if theme == "modern":
        _add_title_bar(slide, colors)

    heading_top = 0.85 if theme == "modern" else 0.5
    _add_text_box(
        slide,
        left_in=0.75, top_in=heading_top,
        width_in=SLIDE_WIDTH_IN - 1.5, height_in=0.9,
        text=heading or "Section",
        font_size_pt=40,
        color_rgb=colors["text"],
        bold=True,
        align="left",
    )

    # Small accent underline below the heading (light/bold themes)
    if theme in ("light", "bold"):
        from pptx.util import Inches
        from pptx.enum.shapes import MSO_SHAPE
        underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.75), Inches(heading_top + 0.95),
            Inches(1.5), Inches(0.06),
        )
        underline.fill.solid()
        underline.fill.fore_color.rgb = _rgb(colors["accent"])
        underline.line.fill.background()

    bullets_top = heading_top + 1.2
    bullets_height = SLIDE_HEIGHT_IN - bullets_top - 0.5

    # Pick body font size — shrink slightly if there are a lot of bullets,
    # but never below the 24pt minimum from the spec.
    n = len(bullets or [])
    if n <= 3:
        body_pt = 28
    elif n <= 5:
        body_pt = 26
    else:
        body_pt = 24

    _add_bullet_block(
        slide,
        left_in=0.95, top_in=bullets_top,
        width_in=SLIDE_WIDTH_IN - 1.9, height_in=bullets_height,
        bullets=bullets or ["(no content)"],
        font_size_pt=body_pt,
        color_rgb=colors["text"],
        accent_rgb=colors["accent2"],
    )


def _add_contact_slide(prs, layout, business_info: dict, *,
                       theme: str) -> None:
    """Final 'thank you' + contact details slide."""
    colors = _theme_colors(theme, slide_index=99)
    slide = prs.slides.add_slide(layout)

    _paint_background(slide, colors["bg"])

    if theme == "modern":
        _add_orb_accent(slide, colors)

    name = (business_info.get("name") or "").strip()
    contact = business_info.get("contact") or {}

    headline = "Thank you" if not name else f"Thank you from {name}"

    _add_text_box(
        slide,
        left_in=0.75, top_in=SLIDE_HEIGHT_IN * 0.28,
        width_in=SLIDE_WIDTH_IN - 1.5, height_in=1.3,
        text=headline,
        font_size_pt=52,
        color_rgb=colors["text"],
        bold=True,
        align="center",
    )

    # Build contact lines from whatever's populated
    lines = []
    if contact.get("phone"):
        lines.append(f"Phone: {contact['phone']}")
    if contact.get("email"):
        lines.append(f"Email: {contact['email']}")
    if contact.get("website"):
        lines.append(f"Web: {contact['website']}")
    if contact.get("address"):
        lines.append(f"{contact['address']}")

    if lines:
        _add_bullet_block(
            slide,
            left_in=SLIDE_WIDTH_IN / 2 - 3.0,
            top_in=SLIDE_HEIGHT_IN * 0.28 + 1.7,
            width_in=6.0, height_in=2.5,
            bullets=lines,
            font_size_pt=26,
            color_rgb=colors["muted"],
            accent_rgb=colors["accent"],
        )

    _add_text_box(
        slide,
        left_in=SLIDE_WIDTH_IN - 2.3, top_in=SLIDE_HEIGHT_IN - 0.45,
        width_in=2.0, height_in=0.35,
        text="Made with Orbi",
        font_size_pt=10,
        color_rgb=colors["muted"],
        align="right",
    )


# ---------------------------------------------------------------------------
# Utility
# ---------------------------------------------------------------------------


def _slugify(s: str, max_len: int = 40) -> str:
    s = (s or "").lower().strip()
    s = re.sub(r"[^a-z0-9]+", "_", s).strip("_")
    if not s:
        s = "deck"
    return s[:max_len]


# ---------------------------------------------------------------------------
# ROUTE SURFACE — for the orchestrator to wire into orbi.py
# ---------------------------------------------------------------------------
#
# Owner-authed (cookie). config + workspace_dir come from the logged-in owner.
# workspace_dir = modules.workspace.workspace_path(config) — defaults to ~/Orbi/.
#
#   POST /api/owner/pptx/build
#       body:   { "topic":        "Memorial Day sale pitch deck for our restaurant",
#                 "slide_count":  7,                    # optional, default 7
#                 "theme":        "modern" | "light" | "bold"  # optional, default "modern"
#               }
#       impl:
#           import pptx_gen
#           from modules.workspace import workspace_path
#           result = pptx_gen.build_deck(config, topic, slide_count, theme,
#                                        business_info=config.get("business"))
#           path   = pptx_gen.save_deck_to_workspace(result["pptx_bytes"], topic,
#                                                    workspace_path(config))
#           return { "filename":     path.name,
#                    "download_url": "/api/owner/workspace/files/decks/" + path.name,
#                    "slide_count":  result["slide_count"],
#                    "outline":      result["outline"],
#                    "theme":        result["theme"],
#                    "bytes":        result["bytes"] }
#
#   POST /api/owner/pptx/outline
#       body:   { "topic": "...", "slide_count": 7 }
#       impl:
#           outline = pptx_gen.build_outline(config, topic, slide_count)
#           return { "outline": outline }
#       Use case: owner wants to review/edit the outline BEFORE the .pptx is
#       rendered. Fast (no .pptx serialization).
#
#   POST /api/owner/pptx/from_outline
#       body:   { "outline": { title, subtitle, slides: [...] },
#                 "theme":   "modern" | "light" | "bold" }
#       impl:
#           pptx_bytes = pptx_gen.render_deck(outline, theme=theme,
#                                              business_info=config.get("business"))
#           path = pptx_gen.save_deck_to_workspace(pptx_bytes,
#                                                  outline.get("title","deck"),
#                                                  workspace_path(config))
#           return { "filename":     path.name,
#                    "download_url": "/api/owner/workspace/files/decks/" + path.name,
#                    "bytes":        len(pptx_bytes) }
#       Use case: render the owner's edited outline. Skips the LLM entirely.
#
# Errors:
#   400 — empty topic, empty outline, or outline missing `slides`
#   500 — only if python-pptx itself fails (should never happen on prod)
#
# ---------------------------------------------------------------------------
